from __future__ import annotations

import json
import math
import re
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "ppt_assets"
PPT_PATH = ROOT / "三维重建与高斯泼溅答辩.pptx"

WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

NAVY = RGBColor(20, 32, 48)
INK = RGBColor(28, 35, 46)
MUTED = RGBColor(92, 104, 118)
BLUE = RGBColor(38, 120, 220)
TEAL = RGBColor(35, 166, 150)
ORANGE = RGBColor(236, 132, 54)
RED = RGBColor(210, 70, 70)
GREEN = RGBColor(72, 156, 92)
PALE = RGBColor(245, 248, 252)
WHITE = RGBColor(255, 255, 255)


def natural_key(path: Path):
    parts = re.split(r"(\d+)", path.name)
    return [int(p) if p.isdigit() else p.lower() for p in parts]


def load_metrics(name: str) -> dict:
    return json.loads((ROOT / "outputs" / name / "metrics.json").read_text(encoding="utf-8"))


def load_gaussians(name: str) -> dict:
    return json.loads((ROOT / "outputs" / name / "gaussians.json").read_text(encoding="utf-8"))


def read_ply(path: Path, max_points: int = 6000) -> tuple[np.ndarray, np.ndarray]:
    lines = path.read_text(encoding="utf-8").splitlines()
    count = 0
    start = 0
    for i, line in enumerate(lines):
        if line.startswith("element vertex"):
            count = int(line.split()[-1])
        if line == "end_header":
            start = i + 1
            break
    rows = []
    colors = []
    for line in lines[start : start + count]:
        vals = line.split()
        if len(vals) < 6:
            continue
        rows.append([float(vals[0]), float(vals[1]), float(vals[2])])
        colors.append([int(vals[3]) / 255, int(vals[4]) / 255, int(vals[5]) / 255])
    pts = np.asarray(rows, dtype=np.float32)
    cols = np.asarray(colors, dtype=np.float32)
    if len(pts) > max_points:
        idx = np.linspace(0, len(pts) - 1, max_points, dtype=int)
        pts = pts[idx]
        cols = cols[idx]
    return pts, cols


def normalize_points(points: np.ndarray) -> np.ndarray:
    if len(points) == 0:
        return points
    center = np.median(points, axis=0)
    centered = points - center
    scale = np.percentile(np.linalg.norm(centered, axis=1), 90)
    return centered / max(float(scale), 1e-6)


def make_input_collage() -> Path:
    out = ASSET_DIR / "input_multiview_collage.png"
    paths = sorted((ROOT / "大作业数据" / "数据1-人体").glob("rgb_*.png"), key=natural_key)[:8]
    frames = [Image.open(p).convert("RGB").resize((260, 260), Image.Resampling.LANCZOS) for p in paths]
    canvas = Image.new("RGB", (4 * 260, 2 * 260), (245, 248, 252))
    for i, im in enumerate(frames):
        x = (i % 4) * 260
        y = (i // 4) * 260
        canvas.paste(im, (x, y))
    draw = ImageDraw.Draw(canvas)
    for i, p in enumerate(paths):
        x = (i % 4) * 260 + 8
        y = (i // 4) * 260 + 8
        draw.rounded_rectangle((x, y, x + 86, y + 24), radius=4, fill=(0, 0, 0))
        draw.text((x + 8, y + 5), p.stem, fill=(255, 255, 255))
    canvas.save(out)
    return out


def make_pointcloud_preview() -> Path:
    out = ASSET_DIR / "pointcloud_preview.png"
    specs = [
        ("VGGT initial", ROOT / "outputs" / "vggt_scene" / "vggt_initial_point_cloud.ply"),
        ("BA output", ROOT / "outputs" / "vggt_scene" / "point_cloud.ply"),
        ("Gaussian params", None),
    ]
    fig = plt.figure(figsize=(12, 4), dpi=180)
    for i, (title, path) in enumerate(specs, start=1):
        ax = fig.add_subplot(1, 3, i, projection="3d")
        if path is None:
            g = load_gaussians("gsplat_data1")
            pts = np.asarray(g["points"], dtype=np.float32)
            cols = np.asarray(g["colors"], dtype=np.float32)
        else:
            pts, cols = read_ply(path)
            pts = normalize_points(pts)
        if len(pts):
            ax.scatter(pts[:, 0], pts[:, 2], pts[:, 1], c=cols, s=2 if i < 3 else 12, alpha=0.92)
        ax.set_title(title, fontsize=11, fontweight="bold")
        ax.set_axis_off()
        ax.view_init(elev=18, azim=-60)
        ax.set_box_aspect((1, 1, 1))
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def make_ba_chart() -> Path:
    out = ASSET_DIR / "ba_rmse_chart.png"
    names = ["vggt_data1", "vggt_data2", "vggt_scene", "gsplat_data1", "gsplat_data2"]
    labels = ["VGGT-D1", "VGGT-D2", "VGGT-Scene", "3DGS-D1", "3DGS-D2"]
    initial = []
    final = []
    for name in names:
        m = load_metrics(name)
        initial.append(float(m["ba_initial_rmse"]))
        final.append(float(m["ba_final_rmse"]))
    x = np.arange(len(labels))
    fig, ax = plt.subplots(figsize=(9.4, 4.8), dpi=180)
    ax.bar(x - 0.18, initial, width=0.36, label="Before BA", color="#EC8436")
    ax.bar(x + 0.18, final, width=0.36, label="After BA", color="#2678DC")
    ax.set_xticks(x)
    ax.set_xticklabels(labels, rotation=12, ha="right")
    ax.set_ylabel("Reprojection RMSE (px)")
    ax.grid(axis="y", alpha=0.25)
    ax.legend(frameon=False)
    for xi, a, b in zip(x, initial, final):
        drop = (a - b) / a * 100 if a else 0
        ax.text(xi, max(a, b) + 0.12, f"-{drop:.0f}%", ha="center", fontsize=9, color="#1F2937")
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def make_gaussian_loss_chart() -> Path:
    out = ASSET_DIR / "gaussian_loss_chart.png"
    names = ["gsplat_data1", "gsplat_data2"]
    labels = ["Data1", "Data2"]
    initial = [float(load_metrics(n)["gaussian_initial_loss"]) for n in names]
    final = [float(load_metrics(n)["gaussian_final_loss"]) for n in names]
    x = np.arange(len(labels))
    fig, ax = plt.subplots(figsize=(7.5, 4.4), dpi=180)
    ax.plot(x, initial, "o-", lw=2.5, label="Initial", color="#EC8436")
    ax.plot(x, final, "o-", lw=2.5, label="Optimized", color="#23A696")
    ax.fill_between(x, initial, final, color="#23A696", alpha=0.12)
    ax.set_xticks(x)
    ax.set_xticklabels(labels)
    ax.set_ylabel("Photometric loss")
    ax.grid(axis="y", alpha=0.25)
    ax.legend(frameon=False)
    for xi, a, b in zip(x, initial, final):
        ax.text(xi, max(a, b) + 0.004, f"{a:.3f}->{b:.3f}", ha="center", fontsize=9)
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def make_mask_chart() -> Path:
    out = ASSET_DIR / "mask_ablation.png"
    with_mask = load_metrics("data1")
    no_mask = load_metrics("data1_nomask")
    labels = ["Mask", "No mask"]
    rmse = [with_mask["ba_final_rmse"], no_mask["ba_final_rmse"]]
    points = [with_mask["gaussian_points"], no_mask["gaussian_points"]]
    fig, ax1 = plt.subplots(figsize=(8, 4.4), dpi=180)
    x = np.arange(2)
    ax1.bar(x - 0.16, rmse, width=0.32, color="#2678DC", label="BA RMSE")
    ax1.set_ylabel("BA final RMSE (px)", color="#2678DC")
    ax1.tick_params(axis="y", labelcolor="#2678DC")
    ax2 = ax1.twinx()
    ax2.bar(x + 0.16, points, width=0.32, color="#23A696", label="Gaussian points")
    ax2.set_ylabel("Gaussian points", color="#23A696")
    ax2.tick_params(axis="y", labelcolor="#23A696")
    ax1.set_xticks(x)
    ax1.set_xticklabels(labels)
    ax1.grid(axis="y", alpha=0.2)
    ax1.set_title("Mask-aware ablation")
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def make_viewer_preview() -> Path:
    out = ASSET_DIR / "webgl_viewer_preview.png"
    g = load_gaussians("gsplat_data1")
    pts = np.asarray(g["points"], dtype=np.float32)
    cols = np.asarray(g["colors"], dtype=np.float32)
    radii = np.asarray(g["radii"], dtype=np.float32)
    order = np.argsort(pts[:, 2])
    pts, cols, radii = pts[order], cols[order], radii[order]
    fig, ax = plt.subplots(figsize=(8, 5), dpi=180)
    ax.set_facecolor("#090B0F")
    sizes = np.clip(radii * 4800, 6, 180)
    ax.scatter(pts[:, 0], pts[:, 1], s=sizes, c=cols, alpha=0.72, linewidths=0)
    ax.set_aspect("equal")
    ax.set_axis_off()
    ax.text(0.02, 0.96, "WebGL viewer: drag rotate / wheel zoom", transform=ax.transAxes, color="white", fontsize=11)
    fig.savefig(out, bbox_inches="tight", pad_inches=0, facecolor="#090B0F")
    plt.close(fig)
    return out


def add_bg(slide, color=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, WIDE_W, WIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.34), Inches(12.1), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = INK
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.92), Inches(11.5), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Microsoft YaHei"
        sp.font.size = Pt(12)
        sp.font.color.rgb = MUTED


def add_bullets(slide, items, x, y, w, h, font_size=17, color=INK, level_gap=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.text = text
        p.level = level
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size - level * 1.5)
        p.font.color.rgb = color
        p.space_after = Pt(7 if not level_gap else 5)
    return box


def add_tag(slide, text, x, y, w, color=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    return shape


def add_card(slide, title, body, x, y, w, h, accent=BLUE):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(250, 252, 255)
    card.line.color.rgb = RGBColor(218, 226, 236)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    t = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.16), w - Inches(0.35), Inches(0.35))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = INK
    b = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.58), w - Inches(0.35), h - Inches(0.7))
    tf = b.text_frame
    tf.clear()
    for i, line in enumerate(body):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(12.5)
        p.font.color.rgb = MUTED
        p.space_after = Pt(4)
    return card


def add_image(slide, path: Path, x, y, w, h):
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def add_table(slide, rows, x, y, w, h, font_size=11):
    table = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h).table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Microsoft YaHei"
                p.font.size = Pt(font_size)
                p.font.color.rgb = WHITE if r == 0 else INK
                if r == 0:
                    p.font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if r == 0 else (RGBColor(246, 249, 253) if r % 2 else WHITE)
    return table


def create_ppt(assets: dict[str, Path]):
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H
    blank = prs.slide_layouts[6]

    # 1
    slide = prs.slides.add_slide(blank)
    add_bg(slide, PALE)
    slide.shapes.add_picture(str(assets["collage"]), Inches(7.3), Inches(0.3), width=Inches(5.5), height=Inches(2.75))
    title = slide.shapes.add_textbox(Inches(0.65), Inches(1.3), Inches(6.4), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = "无标定多视角三维重建与 3D Gaussian Splatting"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY
    add_bullets(slide, ["VGGT 相机/点云初始化", "SciPy Bundle Adjustment", "gsplat CUDA Gaussian 优化", "WebGL 实时交互展示"], Inches(0.78), Inches(3.0), Inches(5.7), Inches(2.2), 18)
    add_tag(slide, "答辩重点：任务闭环 + 可运行展示 + 实验分析", Inches(0.78), Inches(6.35), Inches(4.7), TEAL)

    # 2 scoring
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "评分点与本项目对应", "按题目要求组织答辩内容，避免只讲代码不讲得分点")
    rows = [
        ["评分点", "分值", "对应实现 / 展示材料"],
        ["VGGT 相机参数和初步点云", "3", "VGGT 输出 pose_enc、world_points；展示初始 PLY"],
        ["Bundle Adjustment", "4", "SciPy least_squares 联合优化外参和点云；RMSE 图"],
        ["3D Gaussian 优化与交互渲染", "4", "gsplat rasterization + WebGL viewer"],
        ["VGGT 改进方法", "3", "confidence/mask/keyframe/half precision"],
        ["PPT 制作与答辩", "6", "流程、数据、图表、演示脚本完整"],
        ["技术讲解与展示", "2", "命令、输出目录、viewer 演示"],
        ["BA 对 3DGS 影响", "2", "RMSE 与 Gaussian loss 对照"],
        ["改进实验分析 + 未来方向", "2", "mask ablation + 后续研究"],
    ]
    add_table(slide, rows, Inches(0.6), Inches(1.25), Inches(12.1), Inches(5.85), 10.5)

    # 3 data
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "输入数据：多视角图像，无标定参数", "项目只读取 RGB/mask/视频帧，不依赖外部相机标定文件")
    add_image(slide, assets["collage"], Inches(0.75), Inches(1.35), Inches(6.0), Inches(3.0))
    add_card(slide, "数据读取", ["图像目录：rgb_*.png / msk_*.png", "视频：OpenCV 均匀抽帧", "默认全量：不限制图像数与分辨率"], Inches(7.05), Inches(1.25), Inches(5.3), Inches(1.45), BLUE)
    add_card(slide, "无标定处理", ["VGGT 直接估计内外参", "OpenCV fallback 用经验内参", "后续 BA 统一优化外参和点云"], Inches(7.05), Inches(3.05), Inches(5.3), Inches(1.45), TEAL)
    add_card(slide, "答辩展示", ["打开输入拼图说明多视角来源", "强调没有读取 COLMAP/camera.txt"], Inches(7.05), Inches(4.85), Inches(5.3), Inches(1.25), ORANGE)

    # 4 pipeline
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "整体流程", "从无标定图像到可交互 3DGS viewer 的闭环")
    steps = [
        ("多视角图像/视频", BLUE),
        ("VGGT\n相机 + 初始点云", TEAL),
        ("SIFT tracks\n三角化桥接", ORANGE),
        ("Bundle Adjustment\n外参 + 点云", GREEN),
        ("gsplat\nGaussian 优化", RED),
        ("WebGL\n实时交互渲染", NAVY),
    ]
    x0, y = Inches(0.45), Inches(2.35)
    for i, (text, color) in enumerate(steps):
        x = x0 + Inches(i * 2.12)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(1.8), Inches(1.08))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        p = shape.text_frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        if i < len(steps) - 1:
            arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, x + Inches(1.76), y + Inches(0.36), Inches(0.42), Inches(0.32))
            arr.fill.solid()
            arr.fill.fore_color.rgb = RGBColor(170, 180, 192)
            arr.line.fill.background()
    add_bullets(slide, ["关键点：VGGT dense point map 不直接提供 BA 所需的跨图同名 2D 观测", "解决：用 VGGT 相机初始化 SIFT tracks，再三角化进入 BA", "优势：既利用 VGGT 的无标定能力，又保留传统 BA 的可解释几何残差"], Inches(0.9), Inches(4.45), Inches(11.7), Inches(1.4), 16)

    # 5 VGGT
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "VGGT：相机参数与初步点云", "`vggt_adapter.py` 调用官方 VGGT，输出 pose_enc/world_points/confidence")
    add_image(slide, assets["pc"], Inches(0.55), Inches(1.35), Inches(7.0), Inches(3.4))
    add_card(slide, "实现证据", ["VGGT.from_pretrained(...)", "pose_encoding_to_extri_intri(...)", "world_points_conf 置信度", "输出 vggt_initial_point_cloud.ply"], Inches(7.85), Inches(1.25), Inches(4.8), Inches(2.05), BLUE)
    add_card(slide, "答辩话术", ["这一步解决“无标定”入口问题", "VGGT 结果既保存为初始点云，也作为后续 track/BA 的相机先验"], Inches(7.85), Inches(3.75), Inches(4.8), Inches(1.65), TEAL)

    # 6 BA implementation
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Bundle Adjustment：外参与点云联合优化", "优化变量、残差和鲁棒 loss 均在项目中显式组织")
    add_card(slide, "优化变量", ["R_i, t_i：除第一帧外的相机外参", "X_j：全部重建三维点", "第一帧固定消除 gauge 自由度"], Inches(0.65), Inches(1.35), Inches(3.75), Inches(2.05), BLUE)
    add_card(slide, "目标函数", ["所有 track 的重投影误差", "soft_l1 鲁棒损失", "SciPy least_squares 求解"], Inches(4.8), Inches(1.35), Inches(3.75), Inches(2.05), TEAL)
    add_card(slide, "输出指标", ["initial_reprojection_rmse", "ba_initial_rmse", "ba_final_rmse", "ba_backend"], Inches(8.95), Inches(1.35), Inches(3.75), Inches(2.05), ORANGE)
    add_bullets(slide, ["默认 `--ba-points 0`：不再人为限制 BA 点数", "旧版 OpenCV PnP + 数值 Gauss-Newton 保留为 fallback", "答辩时重点讲：残差如何由 Observation tracks 生成"], Inches(0.85), Inches(4.3), Inches(11.6), Inches(1.6), 17)

    # 7 BA chart
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "BA 实验结果：重投影误差下降", "RMSE 下降说明相机外参与点云几何一致性提高")
    add_image(slide, assets["ba"], Inches(0.85), Inches(1.25), Inches(11.8), Inches(5.35))

    # 8 Gaussian
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "3D Gaussian 优化：gsplat 可微 rasterization", "CUDA 下优先使用开源 gsplat，优化 Gaussian 参数")
    add_card(slide, "优化参数", ["means：三维位置", "scales / quats：大小与方向", "colors：颜色", "opacity：透明度"], Inches(0.65), Inches(1.25), Inches(3.6), Inches(2.35), BLUE)
    add_card(slide, "损失函数", ["多视角渲染图 vs 输入图像", "mask 区域内光度误差", "scale / opacity 正则"], Inches(4.65), Inches(1.25), Inches(3.6), Inches(2.35), TEAL)
    add_card(slide, "运行后端", ["gaussian_backend=gsplat-rasterization", "CUDA: 4 x RTX 3090", "无 CUDA 自动 fallback 到 PyTorch"], Inches(8.65), Inches(1.25), Inches(3.9), Inches(2.35), ORANGE)
    add_image(slide, assets["loss"], Inches(2.25), Inches(4.05), Inches(8.6), Inches(2.6))

    # 9 viewer
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "实时交互展示：WebGL 3DGS Viewer", "不是旧版 CPU Canvas，当前 viewer 使用 WebGL shader 渲染")
    add_image(slide, assets["viewer"], Inches(0.85), Inches(1.25), Inches(7.0), Inches(4.4))
    add_card(slide, "交互能力", ["拖拽旋转", "滚轮缩放", "浏览器直接打开 viewer.html", "WebGL 点精灵高斯透明混合"], Inches(8.25), Inches(1.45), Inches(4.3), Inches(2.1), BLUE)
    add_card(slide, "演示路径", ["python3 -m http.server 8765", "outputs/gsplat_data1/viewer.html", "outputs/vggt_scene/viewer.html"], Inches(8.25), Inches(4.0), Inches(4.3), Inches(1.65), TEAL)

    # 10 BA affects 3DGS
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "BA 是否影响高斯泼溅效果？", "结论：有影响。更准的几何投影让 Gaussian 光度优化更稳定")
    add_image(slide, assets["ba"], Inches(0.55), Inches(1.2), Inches(6.2), Inches(3.2))
    add_image(slide, assets["loss"], Inches(6.95), Inches(1.2), Inches(5.65), Inches(3.2))
    add_bullets(slide, ["BA 降低重投影误差：同一三维点在多视角中落点更准确", "Gaussian 光度监督依赖相机投影；投影错会造成颜色污染和漂浮点", "实验中 BA RMSE 降低，随后 Gaussian loss 继续下降"], Inches(0.9), Inches(4.95), Inches(11.6), Inches(1.2), 16)

    # 11 improvement methods
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "VGGT 改进方法调研", "目标：提高精度或速度，不只停留在 baseline")
    add_card(slide, "Confidence-aware", ["利用 world_points_conf/depth_conf", "低置信度点过滤或降权", "减少遮挡/边界漂浮点"], Inches(0.65), Inches(1.25), Inches(3.75), Inches(2.0), BLUE)
    add_card(slide, "Mask-aware", ["人体前景 mask 限制特征和 loss", "减少绿幕边界噪声", "复杂背景收益更明显"], Inches(4.8), Inches(1.25), Inches(3.75), Inches(2.0), TEAL)
    add_card(slide, "Coarse-to-fine", ["关键帧先稳定相机", "逐步加入相邻帧", "降低冗余和错误传播"], Inches(8.95), Inches(1.25), Inches(3.75), Inches(2.0), ORANGE)
    add_card(slide, "Speed", ["CUDA / autocast 半精度", "VGGT keyframe 输入", "gsplat rasterizer 加速"], Inches(2.65), Inches(4.05), Inches(3.75), Inches(1.7), GREEN)
    add_card(slide, "Weighted BA", ["置信度、mask 距离、残差联合加权", "让困难区域影响更小"], Inches(6.95), Inches(4.05), Inches(3.75), Inches(1.7), RED)

    # 12 improvement experiment
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "改进方法实验分析：Mask-aware", "mask 约束降低背景/边缘噪声，点数更可控")
    add_image(slide, assets["mask"], Inches(0.9), Inches(1.25), Inches(6.7), Inches(4.1))
    rows = [
        ["设置", "BA 后 RMSE", "Gaussian 点数", "分析"],
        ["使用 mask", "1.339 px", "291", "点数更少，前景约束更强"],
        ["不使用 mask", "1.348 px", "337", "背景/边缘点更多"],
    ]
    add_table(slide, rows, Inches(8.0), Inches(1.55), Inches(4.7), Inches(1.75), 10.2)
    add_bullets(slide, ["当前绿幕背景较干净，所以 RMSE 差异不大", "但 mask 让点云来源更可解释，复杂背景中通常收益更明显"], Inches(8.05), Inches(3.8), Inches(4.45), Inches(1.4), 14)

    # 13 no limits
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "运行规模：默认不做人为裁剪", "根据最新要求，运行时间不再由默认参数控制")
    items = [
        "`--max-images 0`：全部图像/视频帧",
        "`--max-size 0`：原始分辨率",
        "`--max-features 0`：SIFT 特征不设人工上限",
        "`--ba-points 0`：BA 优化全部点",
        "`--vggt-max-points 0`：VGGT 点云不截断",
        "`--gaussian-max-points 0`：Gaussian 点不截断",
        "`--gaussian-size 0` / `--gaussian-views 0`：原始分辨率 + 全部视角",
    ]
    add_bullets(slide, items, Inches(0.9), Inches(1.35), Inches(7.0), Inches(4.9), 15)
    add_card(slide, "答辩说明", ["全量运行耗时主要来自 CPU SfM/BA 前处理", "gsplat 第一次运行会编译 CUDA 扩展", "演示可直接打开已生成 viewer 结果"], Inches(8.25), Inches(1.7), Inches(4.25), Inches(2.25), ORANGE)
    add_card(slide, "已验证", ["vggt 环境 CUDA 可用", "4 x RTX 3090", "gsplat-rasterization 后端跑通"], Inches(8.25), Inches(4.35), Inches(4.25), Inches(1.55), TEAL)

    # 14 demo script
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "答辩展示脚本", "建议现场按 3 分钟技术展示组织")
    add_bullets(slide, [
        "1. 展示输入多视角图像：无相机标定文件",
        "2. 打开 `outputs/vggt_scene/vggt_initial_point_cloud.ply` 说明 VGGT 初始点云",
        "3. 展示 BA RMSE 图：解释重投影残差下降",
        "4. 展示 `metrics.json`：backend = scipy + gsplat-rasterization",
        "5. 浏览器打开 `viewer.html`：拖拽旋转、滚轮缩放",
        "6. 讲 BA 对 3DGS 的影响与 mask-aware 改进实验",
    ], Inches(0.85), Inches(1.25), Inches(7.4), Inches(4.8), 17)
    add_card(slide, "现场命令", ["python3 -m http.server 8765", "http://127.0.0.1:8765/outputs/gsplat_data1/viewer.html"], Inches(8.55), Inches(1.75), Inches(4.0), Inches(1.55), BLUE)
    add_card(slide, "若时间紧", ["优先展示 viewer + RMSE 图 + 改进实验", "代码细节留到问答"], Inches(8.55), Inches(4.05), Inches(4.0), Inches(1.45), TEAL)

    # 15 future
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "未来可开展的研究方向", "从作业系统走向更完整的研究型重建")
    add_card(slide, "完整 3DGS 训练", ["引入 densification / pruning", "评估 PSNR / SSIM / LPIPS", "更接近论文原版训练流程"], Inches(0.65), Inches(1.35), Inches(3.85), Inches(2.1), BLUE)
    add_card(slide, "VGGT + BA 融合", ["confidence-weighted BA", "动态 outlier rejection", "关键帧 coarse-to-fine"], Inches(4.85), Inches(1.35), Inches(3.85), Inches(2.1), TEAL)
    add_card(slide, "动态人体建模", ["SMPL / 骨架先验", "非刚体多视角重建", "减少姿态变化影响"], Inches(9.05), Inches(1.35), Inches(3.85), Inches(2.1), ORANGE)
    add_bullets(slide, ["答辩收尾：当前项目完成从无标定图像到 VGGT/BA/3DGS/WebGL viewer 的完整闭环，后续重点是更强的光度指标和动态场景先验。"], Inches(1.0), Inches(4.65), Inches(11.4), Inches(0.8), 18)

    # 16 closing
    slide = prs.slides.add_slide(blank)
    add_bg(slide, NAVY)
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(11.5), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    p.text = "总结"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    add_bullets(slide, [
        "已完成：VGGT 无标定初始化、BA 几何优化、gsplat Gaussian 优化、WebGL 实时交互展示",
        "已分析：BA 对 Gaussian 的影响、mask-aware 改进实验、未来研究方向",
        "答辩重点：用图表和 viewer 证明系统闭环，而不是只展示代码",
    ], Inches(1.0), Inches(2.8), Inches(11.5), Inches(2.0), 20, WHITE)

    prs.save(PPT_PATH)


def main():
    ASSET_DIR.mkdir(exist_ok=True)
    assets = {
        "collage": make_input_collage(),
        "pc": make_pointcloud_preview(),
        "ba": make_ba_chart(),
        "loss": make_gaussian_loss_chart(),
        "mask": make_mask_chart(),
        "viewer": make_viewer_preview(),
    }
    create_ppt(assets)
    print(PPT_PATH)


if __name__ == "__main__":
    main()
